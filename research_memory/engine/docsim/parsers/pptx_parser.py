"""PPTX에서 슬라이드 단위 텍스트·이미지를 추출합니다."""
from __future__ import annotations

import io
from typing import Optional

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from research_memory.engine.docsim.models.schemas import ParseResult
from research_memory.engine.docsim.parsers.common import build_text_records, maybe_image_record, resolve_image_mins
from research_memory.engine.docsim.parsers.image_parser import bytes_to_image


def _shape_text(shape) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    parts = []
    for para in shape.text_frame.paragraphs:
        line = "".join(run.text or "" for run in para.runs).strip()
        if not line and para.text:
            line = para.text.strip()
        if line:
            parts.append(line)
    return "\n".join(parts)


def _collect_images(shape, collected: list, *, file_name: str, location: str, width_min: int, height_min: int):
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            _collect_images(
                child, collected, file_name=file_name, location=location,
                width_min=width_min, height_min=height_min,
            )
        return
    if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
        return
    try:
        blob = shape.image.blob
    except Exception:  # noqa: BLE001
        return
    pil = bytes_to_image(blob)
    if pil is None:
        return
    width, height = pil.size
    rec = maybe_image_record(
        file_name=file_name,
        location=location,
        image_id=f"{file_name}::{location}::img{len(collected)}",
        image_bytes=blob,
        width=width,
        height=height,
        width_min=width_min,
        height_min=height_min,
    )
    if rec:
        collected.append(rec)


def parse_pptx(
    file_name: str,
    file_bytes: bytes,
    min_sentence_length: int = 8,
    min_image_size: Optional[int] = None,
    min_image_width: Optional[int] = None,
    min_image_height: Optional[int] = None,
) -> ParseResult:
    width_min, height_min = resolve_image_mins(min_image_size, min_image_width, min_image_height)

    try:
        prs = Presentation(io.BytesIO(file_bytes))
    except Exception as exc:  # noqa: BLE001
        return ParseResult(
            file_name=file_name,
            success=False,
            error_message=f"PPTX를 열 수 없습니다: {exc}",
        )

    units: list[tuple[int, str, str]] = []
    images = []

    try:
        for idx, slide in enumerate(prs.slides, start=1):
            location = f"슬라이드 {idx}"
            texts: list[str] = []
            for shape in slide.shapes:
                t = _shape_text(shape)
                if t:
                    texts.append(t)
                _collect_images(
                    shape, images, file_name=file_name, location=location,
                    width_min=width_min, height_min=height_min,
                )
            # 노트
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame
                if notes and notes.text and notes.text.strip():
                    texts.append(notes.text.strip())
            units.append((idx, location, "\n".join(texts)))
    except Exception as exc:  # noqa: BLE001
        return ParseResult(
            file_name=file_name,
            success=False,
            error_message=f"PPTX 추출 실패: {exc}",
        )

    sentences, pages, _ = build_text_records(
        file_name, "pptx", units, min_sentence_length=min_sentence_length
    )
    return ParseResult(
        file_name=file_name,
        success=True,
        sentences=sentences,
        images=images,
        pages=pages,
    )
